from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Iterable, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

EMU_W = 13.333
EMU_H = 7.5


def _rgb(hex_color: str, default: str = "#0f236b") -> RGBColor:
    value = (hex_color or default).lstrip("#")
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _pt(px_value: float | int | str) -> Pt:
    # Scene graphs use source-image pixels; PowerPoint expects points.
    return Pt(float(px_value) * 0.72)


def _coord(bbox: List[float], slide_size: List[float]) -> tuple[float, float, float, float]:
    sw, sh = slide_size
    x, y, w, h = bbox
    return x / sw * EMU_W, y / sh * EMU_H, w / sw * EMU_W, h / sh * EMU_H


def _shape_type(kind: str):
    return {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "chevron": MSO_SHAPE.CHEVRON,
        "hexagon": MSO_SHAPE.HEXAGON,
    }.get(kind, MSO_SHAPE.RECTANGLE)


def _set_fill_and_line(shape, style: Dict[str, Any]):
    fill = style.get("fill")
    if fill in (None, "none", "transparent"):
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill, "#ffffff")
        if "transparency" in style:
            shape.fill.transparency = float(style["transparency"])

    stroke = style.get("stroke", style.get("line", "none"))
    if stroke in (None, "none", "transparent"):
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(stroke, "#d9d9d9")
        shape.line.width = Pt(float(style.get("stroke_width", 1)))


def _add_text(slide, item: Dict[str, Any], slide_size: List[float]):
    ix, iy, iw, ih = _coord(item["bbox"], slide_size)
    box = slide.shapes.add_textbox(Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(float(item.get("margin_left", 0)))
    tf.margin_right = Pt(float(item.get("margin_right", 0)))
    tf.margin_top = Pt(float(item.get("margin_top", 0)))
    tf.margin_bottom = Pt(float(item.get("margin_bottom", 0)))
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(item.get("valign", "top"), MSO_ANCHOR.TOP)

    style = item.get("style", {})
    paragraphs = str(item.get("text", "")).split("\n")
    for idx, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(style.get("align", item.get("align", "left")), PP_ALIGN.LEFT)
        p.space_after = Pt(float(style.get("space_after", 0)))
        p.line_spacing = float(style.get("line_spacing", 1.0))
        run = p.add_run()
        run.text = text
        run.font.name = style.get("font", "PingFang SC")
        run.font.size = _pt(style.get("font_size", 16))
        run.font.bold = bool(style.get("bold", False))
        run.font.italic = bool(style.get("italic", False))
        run.font.color.rgb = _rgb(style.get("color", "#222222"))
    return box


def _add_rich_text(slide, item: Dict[str, Any], slide_size: List[float]):
    ix, iy, iw, ih = _coord(item["bbox"], slide_size)
    box = slide.shapes.add_textbox(Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(float(item.get("margin_left", 0)))
    tf.margin_right = Pt(float(item.get("margin_right", 0)))
    tf.margin_top = Pt(float(item.get("margin_top", 0)))
    tf.margin_bottom = Pt(float(item.get("margin_bottom", 0)))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = float(item.get("line_spacing", 1.2))
    for segment in item.get("segments", []):
        run = p.add_run()
        run.text = segment.get("text", "")
        run.font.name = segment.get("font", "PingFang SC")
        run.font.size = _pt(segment.get("font_size", 18))
        run.font.bold = bool(segment.get("bold", False))
        run.font.color.rgb = _rgb(segment.get("color", "#0f236b"))
    return box


def _add_shape(slide, item: Dict[str, Any], slide_size: List[float]):
    ix, iy, iw, ih = _coord(item["bbox"], slide_size)
    shp = slide.shapes.add_shape(_shape_type(item.get("shape", item.get("type", "rect"))), Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    _set_fill_and_line(shp, item.get("style", {}))
    if "rotation" in item:
        shp.rotation = float(item["rotation"])
    if item.get("text"):
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item["text"]
        style = item.get("text_style", item.get("style", {}))
        run.font.name = style.get("font", "PingFang SC")
        run.font.size = _pt(style.get("font_size", 14))
        run.font.bold = bool(style.get("bold", False))
        run.font.color.rgb = _rgb(style.get("color", "#222222"))
    return shp


def _add_line(slide, item: Dict[str, Any], slide_size: List[float]):
    sw, sh = slide_size
    x1, y1 = item["from"]
    x2, y2 = item["to"]
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1 / sw * EMU_W),
        Inches(y1 / sh * EMU_H),
        Inches(x2 / sw * EMU_W),
        Inches(y2 / sh * EMU_H),
    )
    style = item.get("style", {})
    line.line.color.rgb = _rgb(style.get("stroke", style.get("color", "#cccccc")))
    line.line.width = Pt(float(style.get("stroke_width", 1)))
    return line


def _resolve_asset_path(path_value: str, base_dir: str | None = None) -> str:
    path = Path(path_value)
    if path.is_absolute():
        return str(path)
    if base_dir:
        return str(Path(base_dir) / path)
    return str(path)


def add_slide_from_plan(prs: Presentation, plan: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_size = plan.get("slide_size", [1600, 900])

    for item in sorted(plan.get("placements", []), key=lambda p: p.get("z", 0)):
        typ = item.get("type")
        if typ in {"image", "photo", "icon"}:
            ix, iy, iw, ih = _coord(item["bbox"], slide_size)
            image_path = _resolve_asset_path(item["path"], plan.get("asset_base_dir"))
            slide.shapes.add_picture(image_path, Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))
        elif typ == "text":
            _add_text(slide, item, slide_size)
        elif typ == "rich_text":
            _add_rich_text(slide, item, slide_size)
        elif typ == "line":
            _add_line(slide, item, slide_size)
        elif typ in {"rect", "rounded_rect", "ellipse", "triangle", "right_triangle", "parallelogram", "arrow", "chevron", "hexagon", "shape"}:
            _add_shape(slide, item, slide_size)
    return slide


def build_from_plan(plan: Dict[str, Any], out_path: str):
    prs = Presentation()
    prs.slide_width = Inches(EMU_W)
    prs.slide_height = Inches(EMU_H)
    add_slide_from_plan(prs, plan)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def build_from_plans(plans: List[Dict[str, Any]], out_path: str):
    prs = Presentation()
    prs.slide_width = Inches(EMU_W)
    prs.slide_height = Inches(EMU_H)
    for plan in plans:
        add_slide_from_plan(prs, plan)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def build_from_scene(scene: Dict[str, Any], out_path: str):
    placements = []
    for obj in scene.get("objects", []):
        item = dict(obj)
        if obj.get("type") == "card":
            item["type"] = "rounded_rect"
            item["shape"] = "rounded_rect"
        placements.append(item)
    return build_from_plan({"slide_size": scene.get("slide_size", [1600, 900]), "placements": placements}, out_path)
